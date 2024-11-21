import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


def convert_file_to_pptx(pdf_path: str) -> str:
    pptx_path = os.path.splitext(pdf_path)[0] + ".pptx"
    images = convert_from_path(pdf_path)
    presentation = Presentation()
    for image in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])

        image.save("temp.jpg", "JPEG")
        slide.shapes.add_picture("temp.jpg", Inches(0), Inches(0),
                                 width=presentation.slide_width,
                                 height=presentation.slide_height)

    presentation.save(pptx_path)

    return pptx_path

