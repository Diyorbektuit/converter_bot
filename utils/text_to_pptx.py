from pptx import Presentation
from pptx.util import Inches


def split_text_into_chunks(text, max_length):
    """Matnni ma'lum uzunlikdagi qismlarga bo‘lib beradi."""
    words = text.split()
    chunks = []
    current_chunk = []

    for word in words:
        if len(" ".join(current_chunk + [word])) <= max_length:
            current_chunk.append(word)
        else:
            chunks.append(" ".join(current_chunk))
            current_chunk = [word]

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


def create_ppt_from_text(text, output_filename="output_presentation.pptx", max_length=100):
    """Matnni bo‘lib, slaydlarga joylashtiradi va PowerPoint fayli yaratadi."""
    presentation = Presentation()

    # Matnni qismlarga bo‘lamiz
    text_chunks = split_text_into_chunks(text, max_length)

    # Har bir qismni yangi slaydga joylashtiramiz
    for chunk in text_chunks:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Bo'sh slayd
        text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8.5), Inches(6))
        text_frame = text_box.text_frame
        text_frame.text = chunk

        # So'zlarni o'ralishini yoqish
        text_frame.word_wrap = True

    # Taqdimotni saqlash
    presentation.save(output_filename)
    print(f"{output_filename} nomli taqdimot yaratildi.")


# Matnni kiritamiz
sample_text = """Bu juda uzun matn bo‘lishi mumkin. Siz bu yerga o‘zingizning matningizni kiritishingiz mumkin.
Har bir slaydda taxminan bir xil uzunlikdagi matn bo‘ladi.
Matn juda ko‘p bo‘lsa, u avtomatik ravishda bir nechta slaydlarga bo‘linadi."""

# Taqdimotni yaratish
create_ppt_from_text(sample_text)
